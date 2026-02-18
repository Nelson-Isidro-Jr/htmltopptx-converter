#!/usr/bin/env python3
"""
HTML to PPTX - High Quality Edition with Transitions
Slides match image dimensions exactly, high-quality rendering, fade transitions
"""

from flask import Flask, request, jsonify, send_file, send_from_directory, Response
from flask_cors import CORS
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
import os
import asyncio
import base64
import uuid
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from playwright.async_api import async_playwright
from PIL import Image

app = Flask(__name__)

# Security: max upload size 5MB
app.config['MAX_CONTENT_LENGTH'] = 5 * 1024 * 1024

# CORS: restrict to allowed origin in production
allowed_origin = os.environ.get('ALLOWED_ORIGIN', '*')
CORS(app, origins=[allowed_origin] if allowed_origin != '*' else '*')

# Rate limiting
limiter = Limiter(get_remote_address, app=app, default_limits=["100 per minute"])


@app.after_request
def add_security_headers(response):
    response.headers['X-Content-Type-Options'] = 'nosniff'
    response.headers['X-Frame-Options'] = 'SAMEORIGIN'
    response.headers['X-XSS-Protection'] = '1; mode=block'
    response.headers['Referrer-Policy'] = 'strict-origin-when-cross-origin'
    response.headers['Permissions-Policy'] = 'camera=(), microphone=(), geolocation=()'
    return response

presentations = {}

# Default slide dimensions (16:9 aspect ratio)
DEFAULT_WIDTH = 10.0  # inches
DEFAULT_HEIGHT = 5.625  # inches
DEFAULT_DPI = 96


def add_fade_transition(slide, duration_ms=1000):
    """Add fade-in transition to a slide using XML manipulation.
    
    This adds a fade transition that plays when the slide appears.
    The slide will only advance on click, not automatically.
    Duration is in milliseconds (1000ms = 1 second).
    """
    # Create the transition XML
    # p:transition element with fade effect
    # advClick="1" means advance on click (not automatic)
    transition_xml = f'''
    <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
                  spd="med" advClick="1">
        <p:fade />
    </p:transition>
    '''
    
    # Parse the XML
    transition_element = parse_xml(transition_xml)
    
    # Insert transition element into the slide's spTree
    # The transition element should be a direct child of the slide element
    slide_element = slide._element
    
    # Check if there's already a transition element and remove it
    existing_transitions = slide_element.findall(qn('p:transition'))
    for et in existing_transitions:
        slide_element.remove(et)
    
    # Insert the new transition element
    # Insert after cSld element if it exists
    csld = slide_element.find(qn('p:cSld'))
    if csld is not None:
        idx = list(slide_element).index(csld)
        slide_element.insert(idx + 1, transition_element)
    else:
        slide_element.insert(0, transition_element)


def add_push_transition(slide, direction="l"):
    """Alternative: Add push transition (from left by default).
    
    Args:
        slide: The slide object
        direction: 'l' (left), 'r' (right), 'u' (up), 'd' (down)
    """
    transition_xml = f'''
    <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
                  spd="med" advClick="1">
        <p:push dir="{direction}" />
    </p:transition>
    '''
    
    transition_element = parse_xml(transition_xml)
    slide_element = slide._element
    
    # Remove existing transitions
    existing_transitions = slide_element.findall(qn('p:transition'))
    for et in existing_transitions:
        slide_element.remove(et)
    
    # Insert after cSld
    csld = slide_element.find(qn('p:cSld'))
    if csld is not None:
        idx = list(slide_element).index(csld)
        slide_element.insert(idx + 1, transition_element)
    else:
        slide_element.insert(0, transition_element)


async def capture_screenshot(html_content, target_width=None, target_height=None):
    """Capture high-quality screenshot matching slide dimensions.
    
    Args:
        html_content: The HTML/CSS content to render
        target_width: Target width in pixels (optional, defaults to 1920)
        target_height: Target height in pixels (optional, defaults to 1080)
    
    Returns:
        Screenshot bytes and the actual dimensions used
    """
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        
        # Use high resolution for quality (2x scale factor)
        viewport_width = target_width or 1920
        viewport_height = target_height or 1080
        
        # Use device scale factor 2 for retina-quality screenshots
        context = await browser.new_context(
            viewport={'width': viewport_width, 'height': viewport_height},
            device_scale_factor=2  # 2x for high quality
        )
        
        page = await context.new_page()
        
        # Wrap HTML to ensure proper sizing - NO flexbox, just fill the viewport
        wrapped_html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        html, body {{ 
            width: 100%; 
            height: 100%; 
            overflow: hidden;
            background-color: transparent;
        }}
        /* Make the slide-container fill the entire viewport */
        .slide-container {{
            width: 100% !important;
            height: 100% !important;
        }}
    </style>
</head>
<body>{html_content}</body>
</html>"""
        
        await page.set_content(wrapped_html)
        await page.wait_for_load_state('networkidle')
        
        # Wait for fonts, images, and effects to render
        await asyncio.sleep(3)
        
        # Take high-quality screenshot at device scale factor 2
        screenshot = await page.screenshot(
            full_page=False,  # Capture viewport only
            type='png'
        )
        
        await browser.close()
        
        # Get actual dimensions of screenshot
        img = Image.open(BytesIO(screenshot))
        actual_width, actual_height = img.size
        
        return screenshot, actual_width, actual_height


def calculate_slide_dimensions(img_width, img_height, max_width_inches=13.333, max_height_inches=7.5):
    """Calculate slide dimensions to match image aspect ratio.
    
    Args:
        img_width: Image width in pixels
        img_height: Image height in pixels
        max_width_inches: Maximum slide width in inches (default 13.333" for 16:9 at 10" width)
        max_height_inches: Maximum slide height in inches
    
    Returns:
        Tuple of (width_inches, height_inches)
    """
    img_ratio = img_width / img_height
    
    # Standard 16:9 ratio is 1.777...
    target_ratio = 16 / 9
    
    # If image is wider than 16:9, fit to max width
    if img_ratio > target_ratio:
        slide_width = max_width_inches
        slide_height = slide_width / img_ratio
    else:
        # Image is taller, fit to max height
        slide_height = max_height_inches
        slide_width = slide_height * img_ratio
    
    return slide_width, slide_height


def create_slide(prs, screenshot, add_transition=True):
    """Create slide with screenshot filling entire slide exactly.
    
    Args:
        prs: Presentation object
        screenshot: Screenshot bytes
        add_transition: Whether to add fade transition
    
    Returns:
        The created slide
    """
    # Open image to get dimensions
    img = Image.open(BytesIO(screenshot))
    img_width, img_height = img.size
    
    # Calculate slide dimensions to match image aspect ratio
    slide_width_inches, slide_height_inches = calculate_slide_dimensions(img_width, img_height)
    
    # Set presentation dimensions
    prs.slide_width = Inches(slide_width_inches)
    prs.slide_height = Inches(slide_height_inches)
    
    # Create blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Resize image to fit slide exactly at 96 DPI
    target_width_px = int(slide_width_inches * DEFAULT_DPI)
    target_height_px = int(slide_height_inches * DEFAULT_DPI)
    
    # Use high-quality Lanczos resampling
    img_resized = img.resize((target_width_px, target_height_px), Image.Resampling.LANCZOS)
    
    # Save to bytes
    img_bytes = BytesIO()
    img_resized.save(img_bytes, format='PNG', optimize=True)
    img_bytes.seek(0)
    
    # Add picture filling entire slide
    slide.shapes.add_picture(
        img_bytes,
        Inches(0),
        Inches(0),
        width=Inches(slide_width_inches),
        height=Inches(slide_height_inches)
    )
    
    # Add fade-in transition
    if add_transition:
        add_fade_transition(slide)
    
    return slide


@app.route('/api/presentations', methods=['POST'])
def create_presentation():
    """Create a new presentation."""
    pres_id = str(uuid.uuid4())
    presentations[pres_id] = {
        'id': pres_id,
        'slides': [],
        'width': DEFAULT_WIDTH,
        'height': DEFAULT_HEIGHT
    }
    return jsonify({'id': pres_id})


@app.route('/api/presentations/<pres_id>/slides', methods=['POST'])
@limiter.limit("30 per minute")
def add_slide(pres_id):
    """Add a slide to the presentation."""
    if pres_id not in presentations:
        return jsonify({'error': 'Not found'}), 404
    
    html_content = request.json.get('html', '')
    
    # Get optional dimension parameters
    target_width = request.json.get('width', 1280)
    target_height = request.json.get('height', 720)
    
    # Capture screenshot
    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)
    screenshot, actual_width, actual_height = loop.run_until_complete(
        capture_screenshot(html_content, target_width, target_height)
    )
    loop.close()
    
    # Create preview thumbnail
    preview = Image.open(BytesIO(screenshot))
    preview.thumbnail((400, 225))
    p_bytes = BytesIO()
    preview.save(p_bytes, format='PNG', optimize=True)
    preview_b64 = base64.b64encode(p_bytes.getvalue()).decode()
    
    slide_id = str(uuid.uuid4())
    presentations[pres_id]['slides'].append({
        'id': slide_id,
        'screenshot': base64.b64encode(screenshot).decode(),
        'preview': f"data:image/png;base64,{preview_b64}",
        'width': actual_width,
        'height': actual_height
    })
    
    return jsonify({
        'id': slide_id,
        'preview': f"data:image/png;base64,{preview_b64}",
        'width': actual_width,
        'height': actual_height
    })


@app.route('/api/presentations/<pres_id>/slides/<slide_id>', methods=['DELETE'])
def delete_slide(pres_id, slide_id):
    """Delete a slide from the presentation."""
    if pres_id not in presentations:
        return jsonify({'error': 'Not found'}), 404
    
    presentations[pres_id]['slides'] = [
        s for s in presentations[pres_id]['slides'] if s['id'] != slide_id
    ]
    return jsonify({'success': True})


@app.route('/api/presentations/<pres_id>/slides/reorder', methods=['POST'])
def reorder_slides(pres_id):
    """Reorder slides in the presentation."""
    if pres_id not in presentations:
        return jsonify({'error': 'Not found'}), 404
    
    slide_ids = request.json.get('slideIds', [])
    slides = presentations[pres_id]['slides']
    slide_map = {s['id']: s for s in slides}
    presentations[pres_id]['slides'] = [
        slide_map[sid] for sid in slide_ids if sid in slide_map
    ]
    return jsonify({'success': True})


@app.route('/api/presentations/<pres_id>/export', methods=['POST'])
@limiter.limit("10 per minute")
def export_presentation(pres_id):
    """Export the presentation as a PPTX file.
    
    All slides will be sized to match their image dimensions exactly,
    with fade-in transitions automatically applied.
    """
    if pres_id not in presentations:
        return jsonify({'error': 'Not found'}), 404
    
    filename = request.json.get('filename', 'presentation')
    pres_data = presentations[pres_id]
    
    # Create presentation
    prs = Presentation()
    
    # We'll update dimensions for each slide in create_slide
    # First slide sets the presentation dimensions
    if pres_data['slides']:
        # Get first slide dimensions to set initial presentation size
        first_screenshot = base64.b64decode(pres_data['slides'][0]['screenshot'])
        img = Image.open(BytesIO(first_screenshot))
        slide_width, slide_height = calculate_slide_dimensions(img.size[0], img.size[1])
        prs.slide_width = Inches(slide_width)
        prs.slide_height = Inches(slide_height)
    
    # Add each slide (dimensions will be adjusted per slide)
    for slide_data in pres_data['slides']:
        screenshot = base64.b64decode(slide_data['screenshot'])
        create_slide(prs, screenshot, add_transition=True)
    
    # Save to bytes
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    return send_file(
        output,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        as_attachment=True,
        download_name=f'{filename}.pptx'
    )


@app.route('/api/export-both', methods=['POST'])
@limiter.limit("10 per minute")
def export_both_presentations():
    """Export both presentations as separate PPTX files.
    
    Returns both presentations as separate files in the response.
    """
    data = request.json
    pres_id1 = data.get('presentationId')
    pres_id2 = data.get('presentationId2')
    filename = data.get('filename', 'presentation')
    
    if pres_id1 not in presentations or pres_id2 not in presentations:
        return jsonify({'error': 'One or both presentations not found'}), 404
    
    pres_data1 = presentations[pres_id1]
    pres_data2 = presentations[pres_id2]
    
    # Create Presentation 1
    prs1 = Presentation()
    if pres_data1['slides']:
        first_screenshot = base64.b64decode(pres_data1['slides'][0]['screenshot'])
        img = Image.open(BytesIO(first_screenshot))
        slide_width, slide_height = calculate_slide_dimensions(img.size[0], img.size[1])
        prs1.slide_width = Inches(slide_width)
        prs1.slide_height = Inches(slide_height)
    
    for slide_data in pres_data1['slides']:
        screenshot = base64.b64decode(slide_data['screenshot'])
        create_slide(prs1, screenshot, add_transition=True)
    
    # Save Presentation 1 to bytes
    output1 = BytesIO()
    prs1.save(output1)
    output1.seek(0)
    pptx1_base64 = base64.b64encode(output1.getvalue()).decode()
    
    # Create Presentation 2
    prs2 = Presentation()
    if pres_data2['slides']:
        first_screenshot = base64.b64decode(pres_data2['slides'][0]['screenshot'])
        img = Image.open(BytesIO(first_screenshot))
        slide_width, slide_height = calculate_slide_dimensions(img.size[0], img.size[1])
        prs2.slide_width = Inches(slide_width)
        prs2.slide_height = Inches(slide_height)
    
    for slide_data in pres_data2['slides']:
        screenshot = base64.b64decode(slide_data['screenshot'])
        create_slide(prs2, screenshot, add_transition=True)
    
    # Save Presentation 2 to bytes
    output2 = BytesIO()
    prs2.save(output2)
    output2.seek(0)
    pptx2_base64 = base64.b64encode(output2.getvalue()).decode()
    
    return jsonify({
        'pptx1': pptx1_base64,
        'pptx2': pptx2_base64,
        'filename1': f'{filename}_1.pptx',
        'filename2': f'{filename}_2.pptx'
    })


@app.route('/api/presentations/<pres_id>', methods=['GET'])
def get_presentation(pres_id):
    """Get presentation details."""
    if pres_id not in presentations:
        return jsonify({'error': 'Not found'}), 404
    
    pres_data = presentations[pres_id]
    return jsonify({
        'id': pres_data['id'],
        'slides': [
            {'id': s['id'], 'preview': s['preview'], 'width': s.get('width'), 'height': s.get('height')}
            for s in pres_data['slides']
        ]
    })


@app.route('/')
def serve_index():
    """Serve the main index.html file."""
    return send_from_directory('.', 'index.html')


DOMAIN = os.environ.get('DOMAIN', 'https://htmltopptx-converter.onrender.com')


@app.route('/robots.txt')
def robots_txt():
    content = f"""User-agent: *
Allow: /

Sitemap: {DOMAIN}/sitemap.xml
"""
    return Response(content, mimetype='text/plain')


@app.route('/sitemap.xml')
def sitemap_xml():
    content = f"""<?xml version="1.0" encoding="UTF-8"?>
<urlset xmlns="http://www.sitemaps.org/schemas/sitemap/0.9">
  <url>
    <loc>{DOMAIN}/</loc>
    <changefreq>weekly</changefreq>
    <priority>1.0</priority>
  </url>
</urlset>
"""
    return Response(content, mimetype='application/xml')


@app.route('/google9f1c56ac4fb41a96.html')
def google_verification():
    return Response('google-site-verification: google9f1c56ac4fb41a96.html', mimetype='text/plain')


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    debug = os.environ.get('FLASK_DEBUG', 'false').lower() == 'true'
    print("HTML to PPTX Converter")
    print(f"Server: http://localhost:{port}")
    app.run(debug=debug, port=port)
